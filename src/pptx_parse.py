"""
PowerPoint to Slideshow Converter Module for Presentator

This module provides functionality to convert PowerPoint (.pptx) files into
Presentator's slideshow format. Uses python-pptx library for parsing PowerPoint
files and extracts text content, images, and formatting information to create
HTML-based slides compatible with the Presentator system.

Features:
- Text extraction with positioning and formatting preservation
- Image extraction and conversion
- Background color detection
- Slide layout analysis
- HTML generation for web display
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io
import re


def clean_text(text):
    """
    Clean and normalize text content from PowerPoint slides.
    
    Removes extra whitespace, normalizes line breaks, and filters out
    control characters to ensure clean text content for HTML display.
    
    Args:
        text (str): Raw text content from PowerPoint slide
        
    Returns:
        str: Cleaned and normalized text content
    """
    if not text:
        return ""
    
    # Remove extra whitespace and normalize line breaks
    text = re.sub(r'\s+', ' ', text.strip())
    
    # Remove any control characters
    text = re.sub(r'[\x00-\x1f\x7f-\x9f]', '', text)
    
    return text


def extract_slide_content(slide):
    """
    Extract text content from a PowerPoint slide with positioning and formatting.
    
    Analyzes slide shapes to extract text content, preserving layout information
    and basic formatting. Handles different shape types and text positioning.
    
    Args:
        slide: PowerPoint slide object from python-pptx
        
    Returns:
        list: List of content dictionaries containing:
            - text: Extracted text content
            - type: Content type ("text", "title", etc.)
            - position: Positioning information (top, left, width, height)
            - formatting: Text formatting attributes
    """
    content_parts = []
    
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            # Extract formatted text instead of plain text
            formatted_text = extract_formatted_text(shape)
            if formatted_text:
                # Get shape position for sorting
                top_position = getattr(shape, 'top', 0)
                content_parts.append({
                    'text': formatted_text,
                    'top': top_position,
                    'shape': shape
                })
    
    # Sort by vertical position (top to bottom)
    content_parts.sort(key=lambda x: x['top'])
    
    return content_parts


def extract_formatted_text(shape):
    """Extract text with formatting from a shape."""
    if not hasattr(shape, 'text_frame'):
        return clean_text(shape.text) if hasattr(shape, 'text') else ""
    
    formatted_parts = []
    
    for paragraph in shape.text_frame.paragraphs:
        para_parts = []
        
        for run in paragraph.runs:
            text = run.text
            if not text.strip():
                continue
            
            # Start with the text
            formatted_text = text
            
            # Apply formatting based on run properties
            if run.font.bold:
                formatted_text = f"<strong>{formatted_text}</strong>"
            
            if run.font.italic:
                formatted_text = f"<em>{formatted_text}</em>"
            
            if run.font.underline:
                formatted_text = f"<u>{formatted_text}</u>"
            
            # Add font size styling if significantly different from default
            try:
                if run.font.size and run.font.size.pt:
                    font_size = run.font.size.pt
                    if font_size > 18:
                        formatted_text = f'<span style="font-size: {font_size}px;">{formatted_text}</span>'
                    elif font_size < 12:
                        formatted_text = f'<span style="font-size: {font_size}px;">{formatted_text}</span>'
            except:
                pass
            
            # Add color if not default black
            try:
                if run.font.color and run.font.color.rgb:
                    rgb = run.font.color.rgb
                    color = f"#{rgb.r:02x}{rgb.g:02x}{rgb.b:02x}"
                    if color != "#000000":  # Not black
                        formatted_text = f'<span style="color: {color};">{formatted_text}</span>'
            except:
                pass
            
            para_parts.append(formatted_text)
        
        if para_parts:
            formatted_parts.append(''.join(para_parts))
    
    return '\n'.join(formatted_parts) if formatted_parts else ""


def extract_slide_images(slide, presentation_name, output_dir, slide_num):
    """Extract all images from a slide."""
    images = []
    image_count = 0
    
    # Create presentation-specific directory for images
    presentation_dir = output_dir / f"{presentation_name}_images"
    presentation_dir.mkdir(exist_ok=True)
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image_count += 1
                # Get image data
                image_stream = io.BytesIO(shape.image.blob)
                image = Image.open(image_stream)
                
                # Save image with presentation name prefix
                image_filename = f"{presentation_name}_slide_{slide_num}_img_{image_count}.png"
                image_path = presentation_dir / image_filename
                image.save(str(image_path), "PNG")
                
                # Store relative path for HTML
                relative_path = f"{presentation_name}_images/{image_filename}"
                
                images.append({
                    "filename": relative_path,
                    "path": str(image_path)
                })
            except Exception as e:
                print(f"Error extracting image from slide {slide_num}: {e}")
        
    return images


def is_title_text(content_part, index):
    """Determine if text is likely a title."""
    text = content_part['text']
    shape = content_part['shape']
    
    # First text element is likely title
    if index == 0:
        return True
    
    # Short text at top is likely title
    plain_text = re.sub(r'<[^>]+>', '', text)  # Remove HTML tags for length check
    if len(plain_text) < 100 and content_part['top'] < 1000000:  # EMU units
        return True
    
    # Check if text contains large font size indicators
    if 'font-size:' in text and ('18px' in text or '20px' in text or '24px' in text):
        return True
    
    # Check font size if available
    try:
        if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
            para = shape.text_frame.paragraphs[0]
            if para.runs:
                font_size = getattr(para.runs[0].font, 'size', None)
                if font_size and font_size.pt > 20:
                    return True
    except:
        pass
    
    return False


def format_slide_html(content_parts, images):
    """Convert slide content to rich HTML format matching existing slideshows."""
    if not content_parts and not images:
        return "<p>Empty slide</p>"
    
    html_parts = []
    
    # Process text content
    for i, content_part in enumerate(content_parts):
        text = content_part['text']
        
        if is_title_text(content_part, i):
            # Format as title - preserve existing formatting
            plain_text = re.sub(r'<[^>]+>', '', text)  # Remove HTML for length check
            if len(plain_text) > 60:
                html_parts.append(f'<h2>{text}</h2>')
            else:
                html_parts.append(f'<h1 align="center">{text}</h1>')
        else:
            # Process content text - preserve formatting
            lines = [line.strip() for line in text.split('\n') if line.strip()]
            
            list_items = []
            regular_paragraphs = []
            
            for line in lines:
                # Check for list indicators (remove HTML tags first for detection)
                plain_line = re.sub(r'<[^>]+>', '', line)
                if (plain_line.startswith(('', '-', '*', '', '')) or 
                    any(plain_line.startswith(f'{i}.') for i in range(1, 20)) or
                    plain_line.startswith(tuple(f'{i})' for i in range(1, 20)))):
                    # Clean list marker but preserve formatting
                    clean_line = re.sub(r'^(<[^>]*>)*[\-*]?\s*\d*[.)]*\s*', '', line)
                    if clean_line:
                        list_items.append(f'<li>{clean_line}</li>')
                else:
                    if line:
                        # Keep the formatted line as is (formatting already applied)
                        regular_paragraphs.append(f'<p>{line}</p>')
            
            # Add regular paragraphs first
            html_parts.extend(regular_paragraphs)
            
            # Add list if we have items
            if list_items:
                html_parts.append('<ul>')
                html_parts.extend(list_items)
                html_parts.append('</ul>')
    
    # Add images with smart layout
    if images:
        if len(images) == 1:
            # Single image - center it
            img_src = f"/slideshows/{images[0]['filename']}"
            img_html = f'<div align="center"><img class="slide-image" src="{img_src}"/></div>'
            html_parts.append(img_html)
        elif len(images) == 2:
            # Two images - side by side
            html_parts.append('<div style="display: flex; justify-content: center; align-items: flex-start; gap: 10px; margin: 15px 0; flex-wrap: wrap;">')
            for img in images:
                img_src = f"/slideshows/{img['filename']}"
                img_html = f'<div style="flex: 1; max-width: 48%; text-align: center;"><img src="{img_src}" style="max-width: 100%; max-height: 50vh; height: auto; width: auto; border-radius: 6px; box-shadow: 0 2px 6px rgba(0,0,0,0.15); object-fit: contain;" /></div>'
                html_parts.append(img_html)
            html_parts.append('</div>')
        elif len(images) <= 4:
            # 3-4 images - responsive grid
            html_parts.append('<div style="display: flex; justify-content: center; align-items: flex-start; gap: 8px; margin: 15px 0; flex-wrap: wrap; max-width: 100%;">')
            for img in images:
                img_src = f"/slideshows/{img['filename']}"
                img_html = f'<div style="flex: 1; min-width: 200px; max-width: 48%; text-align: center; margin-bottom: 8px;"><img src="{img_src}" style="max-width: 100%; max-height: 40vh; height: auto; width: auto; border-radius: 6px; box-shadow: 0 2px 6px rgba(0,0,0,0.15); object-fit: contain;" /></div>'
                html_parts.append(img_html)
            html_parts.append('</div>')
        else:
            # Many images - compact grid
            html_parts.append('<div style="display: flex; justify-content: center; align-items: flex-start; gap: 5px; margin: 15px 0; flex-wrap: wrap; max-width: 100%;">')
            for img in images:
                img_src = f"/slideshows/{img['filename']}"
                img_html = f'<div style="flex: 1; min-width: 150px; max-width: 30%; text-align: center; margin-bottom: 5px;"><img src="{img_src}" style="max-width: 100%; max-height: 30vh; height: auto; width: auto; border-radius: 4px; box-shadow: 0 1px 4px rgba(0,0,0,0.1); object-fit: contain;" /></div>'
                html_parts.append(img_html)
            html_parts.append('</div>')
    
    return ''.join(html_parts) if html_parts else '<p align="center">Slide content</p>'


def format_text_emphasis(text):
    """Add basic text formatting like bold and italic."""
    # Simple emphasis detection and formatting
    text = re.sub(r'\*\*(.*?)\*\*', r'<strong>\1</strong>', text)
    text = re.sub(r'\*(.*?)\*', r'<em>\1</em>', text)
    text = re.sub(r'_(.*?)_', r'<u>\1</u>', text)
    return text


def extract_slide_background(slide):
    """Try to extract background color from slide."""
    try:
        # Try to get background fill
        if hasattr(slide, 'background') and hasattr(slide.background, 'fill'):
            fill = slide.background.fill
            if hasattr(fill, 'fore_color') and hasattr(fill.fore_color, 'rgb'):
                rgb = fill.fore_color.rgb
                return f"#{rgb.r:02x}{rgb.g:02x}{rgb.b:02x}"
        
        # Generate varied background colors for visual appeal
        colors = ["#f0f8ff", "#fff5ee", "#f5fffa", "#ffefd5", "#f0fff0", "#e6e6fa"]
        return colors[hash(str(slide)) % len(colors)]
    except:
        return "#f8f9fa"


def convert_pptx_to_slideshow_free(pptx_path, output_name=None):
    """
    Convert PowerPoint presentation to Presentator slideshow format.
    
    Main conversion function that processes a PowerPoint file and converts
    it to the Presentator slideshow format with HTML content and extracted images.
    
    Args:
        pptx_path (str or Path): Path to the PowerPoint (.pptx) file
        output_name (str, optional): Name for the output slideshow. If not provided,
            uses the filename without extension
            
    Returns:
        dict: Slideshow data dictionary containing:
            - name: Slideshow name
            - slides: List of slide dictionaries with HTML content
            - metadata: Additional information about the conversion
            
    Raises:
        FileNotFoundError: If the PowerPoint file doesn't exist
        Exception: For various PowerPoint parsing or conversion errors
    """
    pptx_file = Path(pptx_path)
    
    if not pptx_file.exists():
        raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")
    
    # Load presentation
    presentation = Presentation(str(pptx_file))
    
    # Generate output name and clean it for file system
    if not output_name:
        output_name = pptx_file.stem
    
    # Clean name for file system (remove invalid characters)
    clean_name = "".join(c for c in output_name if c.isalnum() or c in (' ', '-', '_')).strip()
    presentation_name = clean_name.replace(' ', '_')
    
    # Create output directory for images
    slideshows_dir = Path("slideshows")
    slideshows_dir.mkdir(exist_ok=True)
    
    slides_data = []
    
    for slide_index, slide in enumerate(presentation.slides):
        slide_num = slide_index + 1
        print(f"Processing slide {slide_num}...")
        
        # Extract images from slide
        images = extract_slide_images(slide, presentation_name, slideshows_dir, slide_num)
        
        # Extract content from slide
        content_parts = extract_slide_content(slide)
        
        # Format as rich HTML
        html_content = format_slide_html(content_parts, images)
        
        # Get background color
        background_color = extract_slide_background(slide)
        
        # Create slide data
        slide_data = {
            "html": html_content,
            "duration": 5000,  # 5 seconds default
            "bgColor": background_color
        }
        
        slides_data.append(slide_data)
        
        # Debug: print first few slides content
        if slide_index < 2:
            print(f"  Text parts: {len(content_parts)}")
            print(f"  Images: {len(images)}")
            print(f"  Background: {background_color}")
    
    # Create slideshow structure
    slideshow_data = {
        "name": output_name,
        "timestamp": "2025-08-01T12:00:00Z",
        "slides": slides_data
    }
    
    return slideshow_data


def save_converted_slideshow_free(slideshow_data, filename=None):
    """Save converted slideshow to JSON file."""
    slideshows_dir = Path("slideshows")
    slideshows_dir.mkdir(exist_ok=True)
    
    if not filename:
        name = slideshow_data.get('name', 'Converted Slideshow')
        # Clean filename
        clean_name = "".join(c for c in name if c.isalnum() or c in (' ', '-', '_')).rstrip()
        filename = f"{clean_name}_editor.json"
    
    if not filename.endswith('_editor.json'):
        filename = filename.replace('.json', '_editor.json')
    
    filepath = slideshows_dir / filename
    
    with open(filepath, 'w', encoding='utf-8') as f:
        json.dump(slideshow_data, f, indent=2, ensure_ascii=False)
    
    return str(filepath)


def cleanup_presentation_files(presentation_name):
    """Remove all files associated with a presentation."""
    slideshows_dir = Path("slideshows")
    if not slideshows_dir.exists():
        return
    
    # Clean presentation name for file system
    clean_name = "".join(c for c in presentation_name if c.isalnum() or c in (' ', '-', '_')).strip()
    file_prefix = clean_name.replace(' ', '_')
    
    # Remove JSON file
    json_file = slideshows_dir / f"{file_prefix}_editor.json"
    if json_file.exists():
        json_file.unlink()
        print(f"Removed: {json_file}")
    
    # Remove image directory
    image_dir = slideshows_dir / f"{file_prefix}_images"
    if image_dir.exists():
        import shutil
        shutil.rmtree(image_dir)
        print(f"Removed image directory: {image_dir}")
    
    print(f"Cleaned up all files for presentation: {presentation_name}")


def convert_pptx_file_free(pptx_path, output_name=None):
    """
    Main entry point for converting PPTX files to Presentator format.
    
    High-level function that handles the complete conversion workflow from
    PowerPoint file to saved slideshow in the Presentator system.
    
    Args:
        pptx_path (str or Path): Path to the PowerPoint file to convert
        output_name (str, optional): Name for the output slideshow
        
    Returns:
        str or None: Path to the saved slideshow file if successful, None if failed
        
    Note:
        This function handles all aspects of conversion including error handling,
        progress reporting, and file saving. It's the recommended entry point
        for PowerPoint conversion in the Presentator system.
    """
    try:
        print(f"Converting PowerPoint file (free version): {pptx_path}")
        
        # Convert presentation
        slideshow_data = convert_pptx_to_slideshow_free(pptx_path, output_name)
        
        # Save to file
        output_path = save_converted_slideshow_free(slideshow_data)
        
        print(f"Conversion completed successfully!")
        print(f"Output file: {output_path}")
        print(f"Slides converted: {len(slideshow_data['slides'])}")
        
        return output_path
        
    except Exception as e:
        print(f"Error converting PowerPoint file: {e}")
        return None


if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python pptx_parse_free.py <pptx_file> [output_name]")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    output_name = sys.argv[2] if len(sys.argv) > 2 else None
    
    convert_pptx_file_free(pptx_file, output_name)
